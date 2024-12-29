import pandas as pd
import pdfplumber
from pptx import Presentation
import json
import re

class DataIngestion:
    def __init__(self, json_path, csv_path, pdf_path, pptx_path):
        self.json_path = json_path
        self.csv_path = csv_path
        self.pdf_path = pdf_path
        self.pptx_path = pptx_path

    def import_json(self):
        return pd.read_json(self.json_path)

    def import_csv(self):
        return pd.read_csv(self.csv_path)

    def import_pdf(self):
        pdf = pdfplumber.open(self.pdf_path)
        page = pdf.pages[0]
        table = page.extract_table()
        return pd.DataFrame(table[1:], columns=table[0])

    def import_pptx(self):
        pptx = Presentation(self.pptx_path)
        text_data = ""
        for slide in pptx.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text
                    text_data += text + "\n"
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_data = [cell.text for cell in row.cells]
                        text_data += " | ".join(row_data) + "\n"
        return text_data

class DataProcessor:
    @staticmethod
    def process_json(json_df):
        companies_df = pd.DataFrame(columns=['id', 'name', 'industry', 'revenue', 'location'])
        companies_revenue_df = pd.DataFrame(columns=['company_id', 'year', 'quarter', 'revenue'])
        companies_profit_margin_df = pd.DataFrame(columns=['company_id', 'year', 'quarter', 'profit_margin'])
        employees_df = pd.DataFrame(columns=['id', 'name', 'company_id', 'role', 'salary', 'hired_date'])

        for index, row in json_df.iterrows():
            company_data = row['companies']
            companies_df = pd.concat([companies_df, pd.DataFrame([{
                'id': company_data.get('id'),
                'name': company_data.get('name'),
                'industry': company_data.get('industry'),
                'revenue': company_data.get('revenue'),
                'location': company_data.get('location')
            }])], ignore_index=True)

            performance_data = company_data.get('performance', {})

            for quarter in performance_data.keys():
                year, qtr = quarter.split('_')
                companies_revenue_df = pd.concat([companies_revenue_df, pd.DataFrame([{
                    'company_id': company_data.get('id'),
                    'year': year,
                    'quarter': qtr,
                    'revenue': performance_data.get(quarter, {}).get('revenue')
                }])], ignore_index=True)

                companies_profit_margin_df = pd.concat([companies_profit_margin_df, pd.DataFrame([{
                    'company_id': company_data.get('id'),
                    'year': year,
                    'quarter': qtr,
                    'profit_margin': performance_data.get(quarter, {}).get('profit_margin')
                }])], ignore_index=True)

            for employee in company_data.get('employees', []):
                employees_df = pd.concat([employees_df, pd.DataFrame([{
                    'id': employee.get('id'),
                    'name': employee.get('name'),
                    'company_id': company_data.get('id'),
                    'role': employee.get('role'),
                    'salary': employee.get('salary'),
                    'hired_date': employee.get('hired_date')
                }])], ignore_index=True)

        return companies_df, companies_revenue_df, companies_profit_margin_df, employees_df

    @staticmethod
    def parse_text_to_json(text_data):
        def parse_key_highlights(data):
            highlights = {}
            revenue_match = re.search(r"Total Revenue:\s*\$([\d,]+)", data)
            memberships_match = re.search(r"Total Memberships Sold:\s*([\d,]+)", data)
            location_match = re.search(r"Top Location:\s*([\w\s]+)", data)

            if revenue_match:
                highlights["TotalRevenue"] = int(revenue_match.group(1).replace(",", ""))
            if memberships_match:
                highlights["TotalMembershipsSold"] = int(memberships_match.group(1).replace(",", ""))
            if location_match:
                highlights["TopLocation"] = location_match.group(1).strip()

            return highlights

        def parse_table(data, header_row_indicator="Quarterly Metrics"):
            rows = []
            table_started = False
            headers = []

            for line in data.split("\n"):
                if header_row_indicator in line:
                    table_started = True
                    continue
                if table_started and "|" in line:
                    row_data = [entry.strip() for entry in line.split("|")]
                    if not headers:
                        headers = row_data
                    else:
                        row = {headers[i]: int(entry.replace(",", "")) if entry.replace(",", "").isdigit() else entry for i, entry in enumerate(row_data)}
                        rows.append(row)

            return rows

        def parse_revenue_breakdown(data):
            breakdown = {}
            matches = re.findall(r"([\w\s]+):\s*(\d+)%", data)
            for match in matches:
                activity, percentage = match
                breakdown[activity.strip()] = int(percentage)
            return breakdown

        result = {}

        if "Key Highlights" in text_data:
            result["KeyHighlights"] = parse_key_highlights(text_data)

        if "Quarterly Metrics" in text_data:
            result["QuarterlyMetrics"] = parse_table(text_data)

        if "Revenue Breakdown by Activity" in text_data:
            result["RevenueBreakdownByActivity"] = parse_revenue_breakdown(text_data)

        return result

class DataSaver:
    @staticmethod
    def save_as_json(df, path):
        df.to_json(path, orient='records')

    @staticmethod
    def save_text_as_json(data, path):
        with open(path, 'w') as f:
            json.dump(data, f, indent=4)